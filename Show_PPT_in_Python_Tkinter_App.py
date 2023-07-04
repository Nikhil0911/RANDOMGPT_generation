import tkinter as tk
from PIL import Image, ImageTk
from pptx import Presentation

class PowerPointViewer(tk.Tk):
    def __init__(self, presentation_file):
        super().__init__()
        self.presentation = Presentation(presentation_file)
        self.current_slide = 0
        self.total_slides = len(self.presentation.slides)

        self.title("PowerPoint Viewer")
        self.geometry("800x600")
        self.configure(bg="white")

        self.slide_image_label = tk.Label(self, bg="white")
        self.slide_image_label.pack(expand=True, fill="both")

        self.previous_button = tk.Button(self, text="Previous", command=self.show_previous_slide)
        self.previous_button.pack(side="left", padx=20, pady=10)

        self.next_button = tk.Button(self, text="Next", command=self.show_next_slide)
        self.next_button.pack(side="right", padx=20, pady=10)

        self.show_slide()

    def show_slide(self):
        slide = self.presentation.slides[self.current_slide]
        slide_image = self.get_slide_image(slide)
        self.update_slide_image(slide_image)

    def show_next_slide(self):
        self.current_slide = (self.current_slide + 1) % self.total_slides
        self.show_slide()

    def show_previous_slide(self):
        self.current_slide = (self.current_slide - 1) % self.total_slides
        self.show_slide()

    def get_slide_image(self, slide):
        slide_width = 720
        slide_height = 540
        slide_image = Image.new("RGB", (slide_width, slide_height), color="white")

        for shape in slide.shapes:
            if shape.shape_type == 13:  # Shape type 13 represents a picture in python-pptx
                image_bytes = shape.image.blob
                image_stream = io.BytesIO(image_bytes)
                image = Image.open(image_stream)
                image = image.resize((slide_width, slide_height))
                slide_image.paste(image, (0, 0))

        return slide_image

    def update_slide_image(self, slide_image):
        slide_image_tk = ImageTk.PhotoImage(slide_image)
        self.slide_image_label.config(image=slide_image_tk)
        self.slide_image_label.image = slide_image_tk

if __name__ == "__main__":
    import io
    import os

    # Replace 'your_presentation_file.pptx' with the actual PowerPoint file path
    presentation_file = "your_presentation_file.pptx"

    if os.path.exists(presentation_file):
        app = PowerPointViewer(presentation_file)
        app.mainloop()
    else:
        print("Presentation file not found.")
